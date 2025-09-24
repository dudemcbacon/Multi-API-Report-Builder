"""
Enhanced Export Manager for Professional Presentation Formats
Provides advanced export capabilities for charts and dashboards
"""
import logging
import tempfile
from pathlib import Path
from typing import Dict, Any, List, Optional
import polars as pl

logger = logging.getLogger(__name__)

try:
    import plotly.graph_objects as go
    import plotly.io as pio
    from PyQt6.QtWidgets import QMessageBox, QProgressDialog
    from PyQt6.QtCore import QThread, pyqtSignal, Qt
    DEPENDENCIES_AVAILABLE = True
except ImportError:
    logger.warning("Export dependencies not available")
    DEPENDENCIES_AVAILABLE = False
    go = None
    pio = None
    QMessageBox = None
    QProgressDialog = None
    QThread = None
    pyqtSignal = None
    Qt = None

class ExportWorker(QThread):
    """Worker thread for export operations"""
    
    progress_updated = pyqtSignal(int)
    export_completed = pyqtSignal(str)
    export_failed = pyqtSignal(str)
    
    def __init__(self, figures: List[go.Figure], export_config: Dict[str, Any]):
        super().__init__()
        self.figures = figures
        self.export_config = export_config
    
    def run(self):
        """Run the export operation"""
        try:
            export_type = self.export_config.get('type', 'html')
            output_path = self.export_config.get('output_path')
            
            if export_type == 'powerpoint':
                self._export_to_powerpoint()
            elif export_type == 'pdf_report':
                self._export_to_pdf_report()
            elif export_type == 'high_res_images':
                self._export_to_high_res_images()
            else:
                self._export_basic(export_type, output_path)
                
        except Exception as e:
            logger.error(f"Export failed: {e}")
            self.export_failed.emit(str(e))
    
    def _export_to_powerpoint(self):
        """Export charts to PowerPoint presentation"""
        try:
            # This would require python-pptx library
            from pptx import Presentation
            from pptx.util import Inches
            import io
            
            prs = Presentation()
            
            for i, fig in enumerate(self.figures):
                # Add slide
                slide_layout = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                title = slide.shapes.title
                title.text = fig.layout.title.text or f"Chart {i+1}"
                
                # Convert plot to image
                img_bytes = fig.to_image(format="png", width=800, height=600)
                img_stream = io.BytesIO(img_bytes)
                
                # Add image to slide
                left = Inches(1)
                top = Inches(1.5)
                height = Inches(5.5)
                slide.shapes.add_picture(img_stream, left, top, height=height)
                
                self.progress_updated.emit(int((i + 1) / len(self.figures) * 100))
            
            # Save presentation
            output_path = self.export_config['output_path']
            prs.save(output_path)
            self.export_completed.emit(output_path)
            
        except ImportError:
            self.export_failed.emit("python-pptx library not installed. Install with: pip install python-pptx")
        except Exception as e:
            self.export_failed.emit(f"PowerPoint export failed: {str(e)}")
    
    def _export_to_pdf_report(self):
        """Export charts as a multi-page PDF report"""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Image, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib.units import inch
            import io
            
            output_path = self.export_config['output_path']
            doc = SimpleDocTemplate(output_path, pagesize=letter)
            
            # Build content
            story = []
            styles = getSampleStyleSheet()
            
            # Add title page
            title = Paragraph("Data Visualization Report", styles['Title'])
            story.append(title)
            story.append(Spacer(1, 0.5*inch))
            
            for i, fig in enumerate(self.figures):
                # Add chart title
                chart_title = fig.layout.title.text or f"Chart {i+1}"
                story.append(Paragraph(chart_title, styles['Heading1']))
                story.append(Spacer(1, 0.2*inch))
                
                # Convert plot to image
                img_bytes = fig.to_image(format="png", width=600, height=400)
                img_stream = io.BytesIO(img_bytes)
                
                # Add image
                img = Image(img_stream, width=6*inch, height=4*inch)
                story.append(img)
                story.append(Spacer(1, 0.5*inch))
                
                self.progress_updated.emit(int((i + 1) / len(self.figures) * 100))
            
            # Build PDF
            doc.build(story)
            self.export_completed.emit(output_path)
            
        except ImportError:
            self.export_failed.emit("reportlab library not installed. Install with: pip install reportlab")
        except Exception as e:
            self.export_failed.emit(f"PDF export failed: {str(e)}")
    
    def _export_to_high_res_images(self):
        """Export charts as high-resolution images"""
        try:
            output_dir = Path(self.export_config['output_path'])
            output_dir.mkdir(exist_ok=True)
            
            for i, fig in enumerate(self.figures):
                # Generate filename
                chart_title = fig.layout.title.text or f"chart_{i+1}"
                # Clean filename
                safe_title = "".join(c for c in chart_title if c.isalnum() or c in (' ', '-', '_')).rstrip()
                filename = f"{safe_title}.png"
                filepath = output_dir / filename
                
                # Export high-res image
                fig.write_image(
                    str(filepath),
                    format='png',
                    width=1920,  # High resolution
                    height=1080,
                    scale=2  # For extra sharpness
                )
                
                self.progress_updated.emit(int((i + 1) / len(self.figures) * 100))
            
            self.export_completed.emit(str(output_dir))
            
        except Exception as e:
            self.export_failed.emit(f"High-res image export failed: {str(e)}")
    
    def _export_basic(self, export_type: str, output_path: str):
        """Export single chart in basic format"""
        if len(self.figures) != 1:
            raise ValueError("Basic export only supports single charts")
        
        fig = self.figures[0]
        
        if export_type == 'html':
            fig.write_html(output_path)
        elif export_type == 'png':
            fig.write_image(output_path, format='png')
        elif export_type == 'pdf':
            fig.write_image(output_path, format='pdf')
        elif export_type == 'svg':
            fig.write_image(output_path, format='svg')
        else:
            raise ValueError(f"Unsupported export type: {export_type}")
        
        self.progress_updated.emit(100)
        self.export_completed.emit(output_path)

class ExportManager:
    """Enhanced export manager for charts and dashboards"""
    
    def __init__(self, parent=None):
        self.parent = parent
        self.export_worker = None
        self.progress_dialog = None
    
    def export_chart(self, figure: go.Figure, export_type: str = 'html', output_path: str = None):
        """Export a single chart"""
        if not DEPENDENCIES_AVAILABLE:
            logger.error("Export dependencies not available")
            return False
        
        try:
            export_config = {
                'type': export_type,
                'output_path': output_path,
                'branding': self._get_branding_config()
            }
            
            self._start_export([figure], export_config)
            return True
            
        except Exception as e:
            logger.error(f"Export initialization failed: {e}")
            if self.parent:
                QMessageBox.critical(self.parent, "Export Error", f"Failed to start export: {str(e)}")
            return False
    
    def export_dashboard(self, figures: List[go.Figure], export_type: str = 'powerpoint', output_path: str = None):
        """Export multiple charts as a dashboard"""
        if not DEPENDENCIES_AVAILABLE:
            logger.error("Export dependencies not available")
            return False
        
        if not figures:
            logger.error("No figures to export")
            return False
        
        try:
            export_config = {
                'type': export_type,
                'output_path': output_path,
                'branding': self._get_branding_config(),
                'dashboard_title': 'Data Dashboard'
            }
            
            self._start_export(figures, export_config)
            return True
            
        except Exception as e:
            logger.error(f"Dashboard export initialization failed: {e}")
            if self.parent:
                QMessageBox.critical(self.parent, "Export Error", f"Failed to start dashboard export: {str(e)}")
            return False
    
    def _start_export(self, figures: List[go.Figure], export_config: Dict[str, Any]):
        """Start the export process with progress dialog"""
        if not self.parent:
            return
        
        # Create progress dialog
        self.progress_dialog = QProgressDialog("Exporting...", "Cancel", 0, 100, self.parent)
        self.progress_dialog.setWindowModality(Qt.WindowModality.WindowModal)
        self.progress_dialog.setAutoClose(True)
        self.progress_dialog.show()
        
        # Create and start export worker
        self.export_worker = ExportWorker(figures, export_config)
        self.export_worker.progress_updated.connect(self.progress_dialog.setValue)
        self.export_worker.export_completed.connect(self._on_export_completed)
        self.export_worker.export_failed.connect(self._on_export_failed)
        self.progress_dialog.canceled.connect(self.export_worker.terminate)
        
        self.export_worker.start()
    
    def _on_export_completed(self, output_path: str):
        """Handle successful export completion"""
        if self.progress_dialog:
            self.progress_dialog.close()
        
        if self.parent:
            QMessageBox.information(
                self.parent, 
                "Export Complete", 
                f"Successfully exported to:\n{output_path}"
            )
        
        logger.info(f"Export completed: {output_path}")
    
    def _on_export_failed(self, error_message: str):
        """Handle export failure"""
        if self.progress_dialog:
            self.progress_dialog.close()
        
        if self.parent:
            QMessageBox.critical(
                self.parent, 
                "Export Failed", 
                f"Export failed:\n{error_message}"
            )
        
        logger.error(f"Export failed: {error_message}")
    
    def _get_branding_config(self) -> Dict[str, Any]:
        """Get branding configuration for exports"""
        return {
            'company_logo': None,  # Path to company logo
            'color_scheme': 'professional',  # Color scheme for branding
            'font_family': 'Arial',
            'footer_text': 'Generated by SalesForce Report Pull'
        }
    
    def get_supported_formats(self) -> List[Dict[str, str]]:
        """Get list of supported export formats"""
        formats = [
            {'name': 'HTML Interactive', 'extension': 'html', 'type': 'html'},
            {'name': 'PNG Image', 'extension': 'png', 'type': 'png'},
            {'name': 'PDF Document', 'extension': 'pdf', 'type': 'pdf'},
            {'name': 'SVG Vector', 'extension': 'svg', 'type': 'svg'},
        ]
        
        # Add advanced formats if dependencies are available
        try:
            import pptx
            formats.append({'name': 'PowerPoint Presentation', 'extension': 'pptx', 'type': 'powerpoint'})
        except ImportError:
            pass
        
        try:
            import reportlab
            formats.append({'name': 'PDF Report', 'extension': 'pdf', 'type': 'pdf_report'})
        except ImportError:
            pass
        
        formats.append({'name': 'High-Res Images', 'extension': 'folder', 'type': 'high_res_images'})
        
        return formats